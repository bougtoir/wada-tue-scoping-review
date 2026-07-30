#!/usr/bin/env python3
"""Prepare clean final figures from previous-version PPTX and create_figures.py.

- Removes figure numbers, captions, and legends from the supplied PPTX files.
- Updates PCOS -> PMOS where needed.
- Exports each figure to PNG/TIFF at 300 dpi.
- Produces individual editable PPTX files and a combined PPTX.
"""

import os
import shutil
import subprocess
import tempfile
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt, Emu

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(BASE_DIR, 'figures')
MANUSCRIPTS_DIR = os.path.join(BASE_DIR, 'manuscripts')
INDIV_DIR = os.path.join(MANUSCRIPTS_DIR, 'individual_figures')

os.makedirs(FIG_DIR, exist_ok=True)
os.makedirs(INDIV_DIR, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Previous-version PPTX sources (new_number -> source file and mapping)
# Note: previous Figure 2 (conceptual) becomes new Figure 4.
#       previous Figure 4 (timeline) becomes new Figure 3.
SOURCES = [
    {
        'new_num': 1,
        'src': '/home/ubuntu/attachments/d4f85261-14c5-4f98-a516-88db5db27370/Figure_1_2.pptx',
        'base': 'fig1_prisma_flowchart',
    },
    {
        'new_num': 3,
        'src': '/home/ubuntu/attachments/88331fe5-f77a-4b0f-bb23-f22c1df92770/Figure_4_2.pptx',
        'base': 'fig3_timeline',
    },
    {
        'new_num': 4,
        'src': '/home/ubuntu/attachments/5d1f8928-6b11-4a1b-aa48-4ddd6ec49458/Figure_2_1.pptx',
        'base': 'fig4_conceptual_framework',
    },
    {
        'new_num': 5,
        'src': '/home/ubuntu/attachments/aedec5f1-92ac-4746-b067-6fd581b42473/Figure_5_2.pptx',
        'base': 'fig5_severity_bar',
    },
]


def pptx_to_png_tiff(pptx_path, out_base):
    """Convert a PPTX slide to PNG/TIFF (300 dpi) using LibreOffice + pdftoppm."""
    with tempfile.TemporaryDirectory() as td:
        stem = os.path.splitext(os.path.basename(pptx_path))[0]
        tmp_pptx = os.path.join(td, f'{stem}.pptx')
        shutil.copy(pptx_path, tmp_pptx)

        pdf_path = os.path.join(td, f'{stem}.pdf')
        subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', td, tmp_pptx],
            check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120,
        )

        png_prefix = os.path.join(td, 'slide')
        subprocess.run(
            ['pdftoppm', '-png', '-r', '300', '-singlefile', pdf_path, png_prefix],
            check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60,
        )
        png_tmp = f'{png_prefix}.png'
        png_out = f'{out_base}.png'
        tiff_out = f'{out_base}.tif'

        shutil.copy(png_tmp, png_out)
        img = Image.open(png_tmp)
        # Save TIFF with LZW compression to keep file sizes reasonable.
        img.save(tiff_out, compression='tiff_lzw', dpi=(300, 300))
        return png_out, tiff_out


def clean_pptx_figures():
    """Generate clean individual PPTX, PNG, and TIFF for the four PPTX sources."""
    for item in SOURCES:
        prs = Presentation(item['src'])
        slide = prs.slides[0]
        to_remove = []

        for sh in slide.shapes:
            # Remove title/caption text boxes
            if sh.has_text_frame:
                txt = sh.text_frame.text
                if txt.startswith('Figure'):
                    to_remove.append(sh)
                else:
                    # Update PCOS -> PMOS in any text
                    if 'PCOS' in txt:
                        for p in sh.text_frame.paragraphs:
                            for r in p.runs:
                                if 'PCOS' in r.text:
                                    r.text = r.text.replace('PCOS', 'PMOS')

            # For the native chart, update categories and keep legend for Figure 5
            if sh.shape_type == 3:
                chart = sh.chart
                series_data = [(s.name, list(s.values)) for s in chart.series]
                cd = ChartData()
                cd.categories = [
                    'Male\nHypogonadism',
                    'ADHD',
                    'PMOS /\nFertility',
                    'Cardiovascular',
                    'Glucocorticoids',
                    'Asthma',
                    'Type 2 Diabetes\n/ GLP-1 RA',
                ]
                for name, vals in series_data:
                    cd.add_series(name, vals)
                chart.replace_data(cd)
                # Figure 5 bar chart should keep its legend; others have no chart.
                chart.has_legend = True
                chart.has_title = False

        for sh in to_remove:
            slide.shapes._spTree.remove(sh._element)

        base = item['base']
        ind_pptx = os.path.join(INDIV_DIR, f'{base}.pptx')
        prs.save(ind_pptx)
        print(f'Saved clean individual PPTX: {ind_pptx}')

        out_base = os.path.join(FIG_DIR, base)
        png, tiff = pptx_to_png_tiff(ind_pptx, out_base)
        print(f'  -> {png}, {tiff}')


def fig2_pptx_from_png():
    """Figure 2 is generated by create_figures.py (clean heatmap). Make an individual PPTX."""
    base = 'fig2_gap_heatmap'
    png_src = os.path.join(FIG_DIR, f'{base}.png')
    tiff_src = os.path.join(FIG_DIR, f'{base}.tif')
    # create_figures.py already wrote these; ensure they exist.
    if not os.path.exists(png_src):
        raise FileNotFoundError(png_src)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img = Image.open(png_src)
    img_w, img_h = img.size
    # Convert pixels at 96 dpi to EMUs, then scale to fit.
    emu_w = int(img_w * 914400 / 96)
    emu_h = int(img_h * 914400 / 96)
    max_w = Inches(12.0)
    max_h = Inches(6.5)
    scale = min(max_w / emu_w, max_h / emu_h, 1.0)
    final_w = int(emu_w * scale)
    final_h = int(emu_h * scale)
    left = int((SLIDE_W - final_w) / 2)
    top = int((SLIDE_H - final_h) / 2)
    slide.shapes.add_picture(png_src, left, top, final_w, final_h)

    ind_pptx = os.path.join(INDIV_DIR, f'{base}.pptx')
    prs.save(ind_pptx)
    print(f'Saved individual PPTX for heatmap: {ind_pptx}')


def create_combined_pptx():
    """Build SCJ_Figures_English.pptx with all five clean figures, no captions."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    order = [
        'fig1_prisma_flowchart.png',
        'fig2_gap_heatmap.png',
        'fig3_timeline.png',
        'fig4_conceptual_framework.png',
        'fig5_severity_bar.png',
    ]

    for fig_file in order:
        fig_path = os.path.join(FIG_DIR, fig_file)
        if not os.path.exists(fig_path):
            print(f'Warning: missing {fig_path}')
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = Image.open(fig_path)
        img_w, img_h = img.size
        emu_w = int(img_w * 914400 / 96)
        emu_h = int(img_h * 914400 / 96)
        max_w = Inches(12.0)
        max_h = Inches(6.5)
        scale = min(max_w / emu_w, max_h / emu_h, 1.0)
        final_w = int(emu_w * scale)
        final_h = int(emu_h * scale)
        left = int((SLIDE_W - final_w) / 2)
        top = int((SLIDE_H - final_h) / 2)
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)

    out = os.path.join(MANUSCRIPTS_DIR, 'SCJ_Figures_English.pptx')
    prs.save(out)
    print(f'Saved combined PPTX: {out}')


if __name__ == '__main__':
    clean_pptx_figures()
    fig2_pptx_from_png()
    create_combined_pptx()
    print('\nAll figure files prepared.')
