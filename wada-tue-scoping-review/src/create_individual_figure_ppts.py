#!/usr/bin/env python3
"""Generate individual editable PPT files for each figure (SCJ submission).

All figures are created using native PowerPoint shapes, text boxes, tables,
and charts so they are fully editable in PowerPoint/Keynote.
Faithfully reproduces the original matplotlib figure designs.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(SCRIPT_DIR, '..', 'manuscripts', 'individual_figures')
os.makedirs(OUT_DIR, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
WADA_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
CLINICAL_GREEN = RGBColor(0x2E, 0x7D, 0x32)
GAP_RED = RGBColor(0xC6, 0x28, 0x28)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)
MEDIUM_GRAY = RGBColor(0x75, 0x75, 0x75)
ORANGE = RGBColor(0xE6, 0x51, 0x00)


def add_title(slide, text, top=Inches(0.15)):
    """Add a centered title text box."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = WADA_BLUE
    return txBox


def add_caption(slide, text, top=Inches(7.0)):
    """Add a centered caption text box."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = DARK_GRAY
    return txBox


def add_box(slide, left, top, width, height, text, fill_rgb, border_rgb,
            font_size=Pt(10), bold=False, text_color=DARK_GRAY):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def add_arrow(slide, start_x, start_y, end_x, end_y, color=None):
    """Add an arrow connector."""
    connector = slide.shapes.add_connector(
        1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color or DARK_GRAY
    connector.line.width = Pt(1.5)
    # Set end arrowhead via XML
    ln = connector.line._ln
    tail_end = ln.makeelement(qn('a:tailEnd'), {})
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    ln.append(tail_end)
    return connector


def add_text(slide, left, top, width, height, text, font_size=Pt(9),
             bold=False, italic=False, color=DARK_GRAY, align=PP_ALIGN.CENTER):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ============================================================
# FIGURE 1: PRISMA-ScR Flow Diagram (faithful reproduction)
# ============================================================
def create_figure_1():
    prs = new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, 'Figure 1. PRISMA-ScR Flow Diagram')

    # Layout constants
    box_w = Inches(2.8)
    box_h = Inches(0.65)
    center_x = Inches(5.2)
    excl_x = Inches(9.5)
    excl_w = Inches(2.3)

    # --- IDENTIFICATION ---
    add_text(slide, Inches(0.3), Inches(0.9), Inches(2.0), Inches(0.3),
             'IDENTIFICATION', Pt(9), True, True, WADA_BLUE, PP_ALIGN.LEFT)

    y1 = Inches(1.1)
    db_box = add_box(slide, Inches(3.0), y1, box_w, box_h,
                     'Records identified through\ndatabase searching (n = 847)',
                     RGBColor(0xBB, 0xDE, 0xFB), WADA_BLUE, Pt(9))
    grey_box = add_box(slide, Inches(7.5), y1, box_w, box_h,
                       'Additional records from grey\nliterature & guidelines (n = 156)',
                       RGBColor(0xBB, 0xDE, 0xFB), WADA_BLUE, Pt(9))

    # Arrows from identification to screening
    add_arrow(slide, Inches(4.4), Inches(1.75), Inches(5.6), Inches(2.05))
    add_arrow(slide, Inches(8.9), Inches(1.75), Inches(7.7), Inches(2.05))

    # --- SCREENING ---
    add_text(slide, Inches(0.3), Inches(2.0), Inches(2.0), Inches(0.3),
             'SCREENING', Pt(9), True, True, WADA_BLUE, PP_ALIGN.LEFT)

    y2 = Inches(2.1)
    dup_box = add_box(slide, Inches(4.2), y2, Inches(4.8), box_h,
                      'Records after duplicates removed (n = 724)',
                      RGBColor(0xC8, 0xE6, 0xC9), CLINICAL_GREEN, Pt(10))

    add_arrow(slide, Inches(6.6), Inches(2.75), Inches(6.6), Inches(3.05))

    y3 = Inches(3.1)
    screen_box = add_box(slide, Inches(4.2), y3, Inches(4.8), box_h,
                         'Records screened by title/abstract (n = 724)',
                         RGBColor(0xC8, 0xE6, 0xC9), CLINICAL_GREEN, Pt(10))
    excl1_box = add_box(slide, excl_x, y3, excl_w, box_h,
                        'Records excluded\n(n = 518)',
                        RGBColor(0xFF, 0xCD, 0xD2), GAP_RED, Pt(9))

    # Arrow to exclusion
    add_arrow(slide, Inches(9.0), Inches(3.4), Inches(9.5), Inches(3.4))
    # Arrow down
    add_arrow(slide, Inches(6.6), Inches(3.75), Inches(6.6), Inches(4.05))

    # --- ELIGIBILITY ---
    add_text(slide, Inches(0.3), Inches(4.0), Inches(2.0), Inches(0.3),
             'ELIGIBILITY', Pt(9), True, True, WADA_BLUE, PP_ALIGN.LEFT)

    y4 = Inches(4.1)
    elig_box = add_box(slide, Inches(4.2), y4, Inches(4.8), box_h,
                       'Full-text articles assessed for eligibility (n = 206)',
                       RGBColor(0xFF, 0xF9, 0xC4), RGBColor(0xF9, 0xA8, 0x25), Pt(10))
    excl2_box = add_box(slide, excl_x, y4, excl_w, box_h,
                        'Full-text excluded\n(n = 138)',
                        RGBColor(0xFF, 0xCD, 0xD2), GAP_RED, Pt(9))

    add_arrow(slide, Inches(9.0), Inches(4.4), Inches(9.5), Inches(4.4))

    # Exclusion reasons box
    add_box(slide, Inches(9.3), Inches(4.8), Inches(3.2), Inches(1.0),
            'Reasons for exclusion:\n- Not athlete-specific (n=42)\n- Outdated guidelines (n=31)\n- Not TUE-related (n=38)\n- Duplicate data (n=27)',
            RGBColor(0xFF, 0xF9, 0xC4), RGBColor(0xBD, 0xBD, 0xBD), Pt(7), False, MEDIUM_GRAY)

    add_arrow(slide, Inches(6.6), Inches(4.75), Inches(6.6), Inches(5.05))

    # --- INCLUDED ---
    add_text(slide, Inches(0.3), Inches(5.0), Inches(2.0), Inches(0.3),
             'INCLUDED', Pt(9), True, True, WADA_BLUE, PP_ALIGN.LEFT)

    y5 = Inches(5.1)
    incl_box = add_box(slide, Inches(4.2), y5, Inches(4.8), box_h,
                       'Sources of evidence included in review (n = 68)',
                       RGBColor(0xC8, 0xE6, 0xC9), RGBColor(0x2E, 0x7D, 0x32), Pt(11), True)

    # Arrows to breakdown
    add_arrow(slide, Inches(5.4), Inches(5.75), Inches(3.5), Inches(6.05))
    add_arrow(slide, Inches(6.6), Inches(5.75), Inches(6.6), Inches(6.05))
    add_arrow(slide, Inches(7.8), Inches(5.75), Inches(9.7), Inches(6.05))

    # Breakdown boxes
    y6 = Inches(6.1)
    bw = Inches(2.5)
    bh = Inches(0.55)
    add_box(slide, Inches(2.3), y6, bw, bh,
            'WADA regulatory\ndocuments (n = 18)',
            RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x15, 0x65, 0xC0), Pt(8))
    add_box(slide, Inches(5.4), y6, bw, bh,
            'Clinical practice\nguidelines (n = 22)',
            RGBColor(0xE8, 0xF5, 0xE9), CLINICAL_GREEN, Pt(8))
    add_box(slide, Inches(8.5), y6, bw, bh,
            'Peer-reviewed\narticles (n = 28)',
            RGBColor(0xFF, 0xF3, 0xE0), ORANGE, Pt(8))

    # Disease area label
    add_text(slide, Inches(3.5), Inches(6.7), Inches(6.0), Inches(0.3),
             'Disease Areas Mapped: Asthma (n=12) | ADHD (n=9) | Diabetes/GLP-1 (n=11) | Hypogonadism (n=10) | GC (n=9) | CV (n=8) | PCOS (n=9)',
             Pt(7), True, False, MEDIUM_GRAY)

    add_caption(slide, 'Figure 1. PRISMA-ScR flow diagram illustrating the source selection process.')

    path = os.path.join(OUT_DIR, 'Figure_1.pptx')
    prs.save(path)
    print(f'Saved: {path}')


# ============================================================
# FIGURE 2: Conceptual Framework (user will edit themselves)
# ============================================================
def create_figure_2():
    """Figure 2 is handled by the user - skip generation."""
    print('Figure 2: Skipped (user will provide their own edit)')


# ============================================================
# FIGURE 3: Gap Risk Matrix (Heatmap style with colored cells)
# ============================================================
def create_figure_3():
    prs = new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, 'Figure 3. Clinical-Competition Gap Risk Matrix\nby Disease Area and Assessment Dimension')

    diseases = [
        'Asthma\n(Beta-2 agonists)',
        'ADHD\n(Stimulants)',
        'Type 2 Diabetes\n(GLP-1 RAs)',
        'Male Hypogonadism\n(Testosterone)',
        'Glucocorticoids\n(Intra-articular)',
        'Cardiovascular\n(Diuretics/BB)',
        'PCOS/Fertility\n(Letrozole/Clomiphene)'
    ]
    dimensions = [
        'First-line Rx\nProhibited',
        'TUE Process\nComplexity',
        'Guideline\nUpdate Lag',
        'Athlete\nImpact',
        'Alternative Rx\nAvailability'
    ]
    data = [
        [2, 2, 2, 3, 4],
        [4, 4, 2, 4, 2],
        [1, 2, 3, 3, 3],
        [4, 4, 3, 4, 1],
        [3, 2, 2, 3, 3],
        [3, 3, 3, 3, 2],
        [4, 3, 2, 3, 1],
    ]
    risk_labels = {0: 'None', 1: 'Low', 2: 'Med', 3: 'High', 4: 'V.High'}

    # Heatmap colors matching original gradient (green -> yellow -> orange -> red)
    severity_colors = {
        0: RGBColor(0xE8, 0xF5, 0xE9),  # Light green
        1: RGBColor(0xC8, 0xE6, 0xC9),  # Green
        2: RGBColor(0xFF, 0xF9, 0xC4),  # Yellow
        3: RGBColor(0xFF, 0xE0, 0xB2),  # Orange
        4: RGBColor(0xEF, 0x53, 0x50),  # Red
    }

    # Grid layout using colored rectangles
    grid_left = Inches(3.0)
    grid_top = Inches(1.2)
    cell_w = Inches(1.7)
    cell_h = Inches(0.7)
    row_label_w = Inches(2.8)

    # Column headers (dimensions)
    for j, dim in enumerate(dimensions):
        x = grid_left + j * cell_w
        add_text(slide, x, grid_top, cell_w, Inches(0.6), dim,
                 Pt(8), True, False, DARK_GRAY)

    # Data rows
    for i, disease in enumerate(diseases):
        y = grid_top + Inches(0.65) + i * cell_h
        # Row label
        add_text(slide, Inches(0.2), y, row_label_w, cell_h, disease,
                 Pt(8), False, False, DARK_GRAY, PP_ALIGN.RIGHT)
        # Cells
        for j, score in enumerate(data[i]):
            x = grid_left + j * cell_w
            color = severity_colors[score]
            text_color = RGBColor(0xFF, 0xFF, 0xFF) if score >= 3 else DARK_GRAY
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, cell_w - Inches(0.05), cell_h - Inches(0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.width = Pt(2)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_top = Pt(0)
            tf.margin_bottom = Pt(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = risk_labels[score]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = text_color

    # Colorbar legend
    legend_y = Inches(6.3)
    legend_labels = ['None', 'Low', 'Medium', 'High', 'Very High']
    legend_colors = [severity_colors[i] for i in range(5)]
    bar_w = Inches(1.5)
    bar_h = Inches(0.35)
    start_x = Inches(3.5)

    add_text(slide, Inches(1.5), legend_y, Inches(2.0), bar_h,
             'Gap Severity:', Pt(9), True, False, DARK_GRAY, PP_ALIGN.RIGHT)

    for k, (lbl, clr) in enumerate(zip(legend_labels, legend_colors)):
        x = start_x + k * bar_w
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, legend_y, bar_w, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = clr
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(1)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lbl
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if k >= 3 else DARK_GRAY

    add_caption(slide,
                'Figure 3. Clinical-competition gap risk matrix by disease area and assessment dimension.')

    path = os.path.join(OUT_DIR, 'Figure_3.pptx')
    prs.save(path)
    print(f'Saved: {path}')


# ============================================================
# FIGURE 4: Timeline (faithful reproduction with all events)
# ============================================================
def create_figure_4():
    prs = new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, 'Figure 4. Timeline of Clinical Guideline Updates vs.\nWADA TUE Regulatory Changes (2018\u20132026)')

    # Central timeline line
    line_y = Inches(3.75)
    line = slide.shapes.add_connector(1, Inches(1.5), line_y, Inches(12.0), line_y)
    line.line.color.rgb = DARK_GRAY
    line.line.width = Pt(2.5)

    # Year markers
    years = list(range(2018, 2027))
    year_spacing = Inches(1.1)
    start_x = Inches(1.7)

    for i, year in enumerate(years):
        x = start_x + i * year_spacing
        # Dot on timeline
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     int(x) - Inches(0.06), int(line_y) - Inches(0.06),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = DARK_GRAY
        dot.line.fill.background()
        # Year label
        add_text(slide, int(x) - Inches(0.3), Inches(3.9), Inches(0.6), Inches(0.25),
                 str(year), Pt(8), True, False, DARK_GRAY)

    # WADA events (above timeline)
    wada_events = [
        (2022, 1, 'GC intra-articular\nprohibited IC'),
        (2023, 10, 'Diabetes TUE\nGuideline v5.1'),
        (2023, 10, 'CV TUE\nGuideline v4.0'),
        (2024, 1, 'GLP-1 RA\nMonitoring'),
        (2025, 1, 'PCOS TUE\nGuideline v2.0'),
        (2025, 12, 'Hypogonadism\nTUE update'),
        (2026, 1, 'Asthma TUE\nv9.3'),
        (2026, 1, 'ADHD TUE\nv8.0'),
    ]

    # Label
    add_text(slide, Inches(0.2), Inches(1.5), Inches(1.3), Inches(0.7),
             'WADA\nRegulatory\nUpdates', Pt(9), True, False, WADA_BLUE, PP_ALIGN.CENTER)

    for i, (year, month, label) in enumerate(wada_events):
        x_pos = start_x + (year - 2018 + month / 12) * year_spacing
        y_level = Inches(1.3) + Inches(0.7) * (i % 3)
        box = add_box(slide, int(x_pos) - Inches(0.7), y_level,
                      Inches(1.5), Inches(0.55), label,
                      RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x15, 0x65, 0xC0), Pt(7))
        # Dashed line to timeline
        conn = slide.shapes.add_connector(1,
                                          int(x_pos), int(y_level) + Inches(0.55),
                                          int(x_pos), int(line_y) - Inches(0.08))
        conn.line.color.rgb = RGBColor(0x15, 0x65, 0xC0)
        conn.line.width = Pt(0.75)
        # Make dashed via XML
        ln = conn.line._ln
        ln.set(qn('a:prstDash') if False else 'prstDash', 'dash')
        prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(prstDash)

    # Clinical events (below timeline)
    clinical_events = [
        (2018, 6, 'Endocrine Soc.\nTestosterone GL'),
        (2022, 1, 'Australian\nADHD GL'),
        (2023, 6, 'PCOS Intl.\nGL revised'),
        (2023, 8, 'TRAVERSE trial\n(TRT safety)'),
        (2023, 10, 'ESC/ESH\nHypertension GL'),
        (2024, 1, 'EAU Sexual/\nReprod. Health'),
        (2024, 6, 'NICE ADHD\nGL update'),
        (2025, 1, 'ADA Diabetes\nStandards 2025'),
        (2025, 5, 'GINA 2025\nAsthma Strategy'),
        (2025, 6, 'ESC Heart\nFailure GL'),
    ]

    # Label
    add_text(slide, Inches(0.2), Inches(4.8), Inches(1.3), Inches(0.7),
             'Clinical\nGuideline\nUpdates', Pt(9), True, False, CLINICAL_GREEN, PP_ALIGN.CENTER)

    for i, (year, month, label) in enumerate(clinical_events):
        x_pos = start_x + (year - 2018 + month / 12) * year_spacing
        y_level = Inches(4.5) + Inches(0.7) * (i % 3)
        box = add_box(slide, int(x_pos) - Inches(0.7), y_level,
                      Inches(1.5), Inches(0.55), label,
                      RGBColor(0xE8, 0xF5, 0xE9), CLINICAL_GREEN, Pt(7))
        # Dashed line to timeline
        conn = slide.shapes.add_connector(1,
                                          int(x_pos), int(line_y) + Inches(0.08),
                                          int(x_pos), int(y_level))
        conn.line.color.rgb = CLINICAL_GREEN
        conn.line.width = Pt(0.75)
        ln = conn.line._ln
        prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(prstDash)

    # GAP arrow and label
    gap_x = Inches(11.5)
    add_arrow(slide, gap_x, Inches(3.5), gap_x, Inches(3.2), GAP_RED)
    add_arrow(slide, gap_x, Inches(4.0), gap_x, Inches(4.3), GAP_RED)
    add_text(slide, Inches(11.2), Inches(3.3), Inches(0.8), Inches(0.4),
             'GAP', Pt(12), True, False, GAP_RED)

    add_caption(slide,
                'Figure 4. Timeline of clinical guideline updates versus WADA TUE regulatory changes (2018\u20132026).')

    path = os.path.join(OUT_DIR, 'Figure_4.pptx')
    prs.save(path)
    print(f'Saved: {path}')


# ============================================================
# FIGURE 5: Bar Chart (editable PowerPoint chart object)
# ============================================================
def create_figure_5():
    prs = new_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, 'Figure 5. Clinical-Competition Gap Severity\nby Disease Area and Assessment Dimension')

    # Chart data
    diseases = [
        'Male\nHypogonadism', 'ADHD', 'PCOS/\nFertility',
        'Cardiovascular', 'Glucocorticoids', 'Asthma',
        'Type 2 Diabetes\n/ GLP-1 RA'
    ]
    first_line = [4, 4, 4, 3, 3, 2, 1]
    tue_complex = [4, 4, 3, 3, 2, 2, 2]
    update_lag = [3, 2, 2, 3, 2, 2, 3]
    athlete_impact = [4, 4, 3, 3, 3, 3, 3]

    chart_data = CategoryChartData()
    chart_data.categories = diseases
    chart_data.add_series('First-line Rx Prohibited', first_line)
    chart_data.add_series('TUE Process Complexity', tue_complex)
    chart_data.add_series('Guideline Update Lag', update_lag)
    chart_data.add_series('Athlete Impact', athlete_impact)

    # Add chart to slide
    chart_left = Inches(0.8)
    chart_top = Inches(1.0)
    chart_width = Inches(11.5)
    chart_height = Inches(5.5)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        chart_left, chart_top, chart_width, chart_height,
        chart_data
    )

    chart = chart_frame.chart

    # Style the chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Color the series
    series_colors = [
        RGBColor(0xEF, 0x53, 0x50),  # Red
        RGBColor(0xFF, 0x98, 0x00),  # Orange
        RGBColor(0xFD, 0xD8, 0x35),  # Yellow
        RGBColor(0x42, 0xA5, 0xF5),  # Blue
    ]
    for i, color in enumerate(series_colors):
        series = chart.series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    # Y-axis settings
    value_axis = chart.value_axis
    value_axis.maximum_scale = 4.5
    value_axis.minimum_scale = 0
    value_axis.major_unit = 1

    add_caption(slide,
                'Figure 5. Clinical-competition gap severity by disease area and assessment dimension.\n'
                'Severity: None (0), Low (1), Medium (2), High (3), Very High (4).',
                top=Inches(6.7))

    path = os.path.join(OUT_DIR, 'Figure_5.pptx')
    prs.save(path)
    print(f'Saved: {path}')


# ============================================================
if __name__ == '__main__':
    create_figure_1()
    create_figure_2()
    create_figure_3()
    create_figure_4()
    create_figure_5()
    print(f'\nAll individual editable figure PPT files saved to {OUT_DIR}')
