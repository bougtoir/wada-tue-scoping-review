#!/usr/bin/env python3
"""Generate editable PPTX with all figures (1 per slide), EN captions."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(SCRIPT_DIR, '..', 'figures')
OUT_DIR = os.path.join(SCRIPT_DIR, '..', 'manuscripts')
os.makedirs(OUT_DIR, exist_ok=True)

# Widescreen dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FIGURES_EN = [
    ('fig1_prisma_flowchart.png',
     'Figure 1',
     'PRISMA-ScR flow diagram illustrating the source selection process.'),
    ('fig2_gap_heatmap.png',
     'Figure 2',
     'Clinical-competition gap risk matrix by disease area and assessment dimension.'),
    ('fig3_timeline.png',
     'Figure 3',
     'Timeline of clinical guideline updates versus WADA TUE regulatory changes (2018\u20132026).'),
    ('fig4_conceptual_framework.png',
     'Figure 4',
     'Conceptual framework: structural drivers of clinical-competition divergence.'),
    ('fig5_severity_bar.png',
     'Figure 5',
     'Clinical-competition gap severity by disease area and assessment dimension.'),
]

FIGURES_JP = [
    ('fig1_prisma_flowchart.png',
     '\u56f3 1',
     'PRISMA-ScR\u30d5\u30ed\u30fc\u30c0\u30a4\u30a2\u30b0\u30e9\u30e0\uff1a\u6587\u732e\u9078\u629e\u30d7\u30ed\u30bb\u30b9\u3002'),
    ('fig2_gap_heatmap.png',
     '\u56f3 2',
     '\u81e8\u5e8a\u2500\u7af6\u6280\u30ae\u30e3\u30c3\u30d7\u30ea\u30b9\u30af\u30de\u30c8\u30ea\u30c3\u30af\u30b9\uff1a\u75be\u60a3\u9818\u57df\u00d7\u8a55\u4fa1\u6b21\u5143\u3002'),
    ('fig3_timeline.png',
     '\u56f3 3',
     '\u81e8\u5e8a\u30ac\u30a4\u30c9\u30e9\u30a4\u30f3\u66f4\u65b0\u3068WADA TUE\u898f\u5236\u5909\u66f4\u306e\u30bf\u30a4\u30e0\u30e9\u30a4\u30f3\uff082018\u20132026\uff09\u3002'),
    ('fig4_conceptual_framework.png',
     '\u56f3 4',
     '\u6982\u5ff5\u7684\u30d5\u30ec\u30fc\u30e0\u30ef\u30fc\u30af\uff1a\u81e8\u5e8a\u2500\u7af6\u6280\u4e56\u96e2\u306e\u69cb\u9020\u7684\u99c6\u52d5\u8981\u56e0\u3002'),
    ('fig5_severity_bar.png',
     '\u56f3 5',
     '\u75be\u60a3\u9818\u57df\u00d7\u8a55\u4fa1\u6b21\u5143\u5225\u306e\u81e8\u5e8a\u2500\u7af6\u6280\u30ae\u30e3\u30c3\u30d7\u91cd\u75c7\u5ea6\u3002'),
]


def add_figure_slide(prs, fig_file, title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title at top
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)

    # Image centered
    fig_path = os.path.join(FIG_DIR, fig_file)
    if os.path.exists(fig_path):
        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        
        max_w = Inches(11.0)
        max_h = Inches(5.0)
        
        scale_w = max_w / Emu(int(img_w * 914400 / 96))
        scale_h = max_h / Emu(int(img_h * 914400 / 96))
        scale = min(scale_w, scale_h, 1.0)
        
        final_w = int(img_w * 914400 / 96 * scale)
        final_h = int(img_h * 914400 / 96 * scale)
        
        left = int((SLIDE_W - final_w) / 2)
        top = Inches(1.0)
        
        slide.shapes.add_picture(fig_path, left, top, final_w, final_h)
    else:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'[Image placeholder: {fig_file}]'
        p2.alignment = PP_ALIGN.CENTER

    # Caption at bottom
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = f'{title}. {caption}'
    run3.font.size = Pt(14)
    run3.font.italic = True


def create_pptx(figures, filename):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for fig_file, title, caption in figures:
        add_figure_slide(prs, fig_file, title, caption)

    path = os.path.join(OUT_DIR, filename)
    prs.save(path)
    print(f'Saved: {path}')


if __name__ == '__main__':
    create_pptx(FIGURES_EN, 'SCJ_Figures_English.pptx')
    create_pptx(FIGURES_JP, 'SCJ_Figures_Japanese.pptx')
